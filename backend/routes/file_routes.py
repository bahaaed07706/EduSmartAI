# routes/file_routes.py - نظام رفع الملفات الحقيقي
"""
Real File Upload System:
1. Actual file upload to server
2. Text extraction from PDF/DOCX
3. Content stored for RAG chatbot
"""
from fastapi import APIRouter, Depends, HTTPException, UploadFile, File, Form
from fastapi.responses import FileResponse
from sqlalchemy.orm import Session
from typing import Optional
import shutil
from datetime import datetime
from pathlib import Path
from database import get_db
from auth import get_lecturer, get_current_user
import models

router = APIRouter(prefix="/files", tags=["Files"])

# Upload directory (resolved once so we can validate paths stay inside it)
UPLOAD_DIR = (Path(__file__).parent.parent / "uploads").resolve()


def _resolve_within_uploads(stored_path: str) -> Optional[Path]:
    """Resolve a stored file reference and confirm it stays inside UPLOAD_DIR.

    Stored references come in three historical shapes, and all must resolve
    identically on Windows and POSIX:

    * ``/uploads/submissions/user_5/ab.pdf`` — the URL-style value returned by
      ``POST /files/upload``. The leading slash makes ``Path.is_absolute()``
      True on POSIX but False on Windows, so it must be stripped *before* any
      absolute/relative decision. Without this the download 404s on Linux while
      passing on a Windows dev machine.
    * ``submissions/user_5/ab.pdf`` — a plain relative path.
    * ``C:\\...\\uploads\\course_1\\x.pdf`` — a real absolute filesystem path
      written by older ``upload_course_material`` rows.

    Guards against traversal: any resolved path that escapes the uploads root
    (via ``..`` or an absolute path elsewhere) is rejected. Returns the safe
    Path, or None if the file is missing or outside the root.
    """
    if not stored_path:
        return None
    try:
        # Normalise separators so a Windows-style value resolves on POSIX too.
        text = str(stored_path).replace("\\", "/").strip()
        if not text:
            return None

        # A genuine absolute path already inside the uploads root (legacy
        # `upload_course_material` rows) is used as-is. This must be tried
        # BEFORE prefix-stripping, so a POSIX path such as
        # /srv/app/backend/uploads/course_1/x.pdf is not mangled.
        candidate = None
        probe = Path(text)
        if probe.is_absolute():
            try:
                resolved = probe.resolve()
                resolved.relative_to(UPLOAD_DIR)
                candidate = resolved
            except (ValueError, OSError):
                candidate = None

        if candidate is None:
            # URL-style or relative reference. Strip the "/uploads" prefix and
            # any leading slashes so it is interpreted relative to UPLOAD_DIR
            # rather than the filesystem root — the leading slash is what makes
            # is_absolute() True on POSIX but False on Windows.
            rel = text
            lowered = rel.lower()
            if lowered.startswith("/uploads/"):
                rel = rel[len("/uploads/"):]
            elif lowered.startswith("uploads/"):
                rel = rel[len("uploads/"):]
            rel = rel.lstrip("/")
            if not rel:
                return None
            candidate = (UPLOAD_DIR / rel).resolve()
            candidate.relative_to(UPLOAD_DIR)  # raises ValueError if outside
    except (ValueError, OSError):
        return None
    return candidate if candidate.is_file() else None

def _invalidate_rag_index() -> None:
    """Drop the cached retrieval index after the material corpus changes.

    Imported lazily to avoid a circular import (chatbot_routes imports nothing
    from here, but the router registration order makes a top-level import
    fragile). Failing to invalidate must never break the upload itself.
    """
    try:
        from routes.chatbot_routes import reset_rag_index

        reset_rag_index()
    except Exception as e:  # pragma: no cover - defensive
        print(f"[files] could not invalidate RAG index: {type(e).__name__}")


# Allowed file types
ALLOWED_EXTENSIONS = {".pdf", ".docx", ".doc", ".txt", ".pptx", ".xlsx"}
MAX_FILE_SIZE = 10 * 1024 * 1024  # 10MB


def extract_text_from_file(file_path: Path) -> str:
    """Extract text content from uploaded file for RAG"""
    text = ""
    suffix = file_path.suffix.lower()
    
    try:
        if suffix == ".txt":
            with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
                text = f.read()
        
        elif suffix == ".pdf":
            try:
                import fitz  # PyMuPDF
                doc = fitz.open(str(file_path))
                for page in doc:
                    text += page.get_text()
                doc.close()
            except ImportError:
                # Fallback: try pdfplumber
                try:
                    import pdfplumber
                    with pdfplumber.open(file_path) as pdf:
                        for page in pdf.pages:
                            page_text = page.extract_text()
                            if page_text:
                                text += page_text + "\n"
                except ImportError:
                    text = "[PDF text extraction not available - install PyMuPDF or pdfplumber]"
        
        elif suffix in {".docx", ".doc"}:
            try:
                import docx
                doc = docx.Document(str(file_path))
                text = "\n".join([para.text for para in doc.paragraphs])
            except ImportError:
                text = "[DOCX text extraction not available - install python-docx]"
        
        elif suffix == ".pptx":
            try:
                from pptx import Presentation
                prs = Presentation(str(file_path))
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "text"):
                            text += shape.text + "\n"
            except ImportError:
                text = "[PPTX text extraction not available - install python-pptx]"
    
    except Exception as e:
        text = f"[Error extracting text: {str(e)}]"
    
    return text.strip()


@router.post("/upload/{course_id}")
async def upload_course_material(
    course_id: int,
    file: UploadFile = File(...),
    title: str = Form(...),
    description: str = Form(""),
    current_user: models.User = Depends(get_lecturer),
    db: Session = Depends(get_db)
):
    """Upload a file as course material (Lecturer only)"""
    
    # Verify lecturer owns this course
    course = db.query(models.Course).filter(
        models.Course.id == course_id,
        models.Course.lecturer_id == current_user.id
    ).first()
    
    if not course:
        raise HTTPException(status_code=403, detail="You don't teach this course")
    
    # Validate file
    if not file.filename:
        raise HTTPException(status_code=400, detail="No file provided")
    
    file_ext = Path(file.filename).suffix.lower()
    if file_ext not in ALLOWED_EXTENSIONS:
        raise HTTPException(
            status_code=400, 
            detail=f"File type not allowed. Allowed: {', '.join(ALLOWED_EXTENSIONS)}"
        )
    
    # Check file size
    file.file.seek(0, 2)  # Seek to end
    file_size = file.file.tell()
    file.file.seek(0)  # Reset
    
    if file_size > MAX_FILE_SIZE:
        raise HTTPException(status_code=400, detail="File too large (max 10MB)")
    
    # Create course folder
    course_folder = UPLOAD_DIR / f"course_{course_id}"
    course_folder.mkdir(parents=True, exist_ok=True)
    
    # Generate unique filename
    timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
    safe_filename = f"{timestamp}_{file.filename}"
    file_path = course_folder / safe_filename
    
    # Save file
    try:
        with open(file_path, "wb") as buffer:
            shutil.copyfileobj(file.file, buffer)
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Failed to save file: {str(e)}")
    
    # Extract text for RAG
    extracted_text = extract_text_from_file(file_path)
    
    # Save to database
    material = models.CourseMaterial(
        course_id=course_id,
        title=title,
        description=description,
        file_url=str(file_path),
        content_text=extracted_text[:10000] if extracted_text else None,  # Limit stored text
        uploaded_by=current_user.id
    )
    db.add(material)
    db.commit()
    db.refresh(material)

    # The retrieval index is process-cached; without this the new material is
    # invisible to the chatbot until the server restarts.
    _invalidate_rag_index()

    return {
        "message": "File uploaded successfully",
        "material_id": material.id,
        "file_name": file.filename,
        "text_extracted": len(extracted_text) > 0,
        "text_length": len(extracted_text)
    }


@router.get("/course/{course_id}")
def get_course_files(
    course_id: int,
    current_user: models.User = Depends(get_current_user),
    db: Session = Depends(get_db)
):
    """Get all files for a course"""
    
    # Check access (actively enrolled student, or the lecturer who owns it).
    if current_user.role == "student":
        enrollment = db.query(models.Enrollment).filter(
            models.Enrollment.student_id == current_user.id,
            models.Enrollment.course_id == course_id,
            # A withdrawn enrollment row still exists; it must not grant access.
            models.Enrollment.status != "withdrawn",
        ).first()
        if not enrollment:
            raise HTTPException(status_code=403, detail="Not enrolled in this course")

    elif current_user.role == "lecturer":
        course = db.query(models.Course).filter(
            models.Course.id == course_id,
            models.Course.lecturer_id == current_user.id
        ).first()
        if not course:
            raise HTTPException(status_code=403, detail="You don't teach this course")

    else:
        # Deny by default: no role falls through to the material list unchecked.
        raise HTTPException(status_code=403, detail="Access denied")

    # Get materials
    materials = db.query(models.CourseMaterial).filter(
        models.CourseMaterial.course_id == course_id
    ).all()
    
    return [
        {
            "id": m.id,
            "title": m.title,
            "description": m.description,
            "file_name": Path(m.file_url).name if m.file_url else None,
            "has_text": bool(m.content_text),
            "created_at": str(m.created_at) if m.created_at else None
        }
        for m in materials
    ]


@router.get("/{material_id}/download")
def download_material(
    material_id: int,
    current_user: models.User = Depends(get_current_user),
    db: Session = Depends(get_db)
):
    """Download a course material file (enrolled students or owning lecturer only)."""

    material = db.query(models.CourseMaterial).filter(
        models.CourseMaterial.id == material_id
    ).first()
    if not material:
        raise HTTPException(status_code=404, detail="Material not found")

    # Access control: student must be enrolled; lecturer must own the course.
    if current_user.role == "student":
        enrollment = db.query(models.Enrollment).filter(
            models.Enrollment.student_id == current_user.id,
            models.Enrollment.course_id == material.course_id,
            models.Enrollment.status != "withdrawn",
        ).first()
        if not enrollment:
            raise HTTPException(status_code=403, detail="Not enrolled in this course")
    elif current_user.role == "lecturer":
        course = db.query(models.Course).filter(
            models.Course.id == material.course_id,
            models.Course.lecturer_id == current_user.id
        ).first()
        if not course:
            raise HTTPException(status_code=403, detail="You don't teach this course")
    else:
        raise HTTPException(status_code=403, detail="Access denied")

    safe_path = _resolve_within_uploads(material.file_url)
    if not safe_path:
        raise HTTPException(status_code=404, detail="File not found on server")

    return FileResponse(
        path=str(safe_path),
        filename=safe_path.name,
        media_type="application/octet-stream"
    )


@router.delete("/{material_id}")
def delete_material(
    material_id: int,
    current_user: models.User = Depends(get_lecturer),
    db: Session = Depends(get_db)
):
    """Delete a material (Lecturer only)"""
    
    material = db.query(models.CourseMaterial).filter(
        models.CourseMaterial.id == material_id
    ).first()
    
    if not material:
        raise HTTPException(status_code=404, detail="Material not found")
    
    # Verify ownership
    course = db.query(models.Course).filter(
        models.Course.id == material.course_id,
        models.Course.lecturer_id == current_user.id
    ).first()
    
    if not course:
        raise HTTPException(status_code=403, detail="You don't own this material")
    
    # Delete file from disk (only if it resolves safely inside the uploads root)
    safe_path = _resolve_within_uploads(material.file_url)
    if safe_path:
        try:
            safe_path.unlink()
        except OSError as e:
            # Do not fail the request if the file is already gone; log for ops.
            print(f"[files] failed to delete file for material {material_id}: {type(e).__name__}")

    # Delete from database
    db.delete(material)
    db.commit()

    # Drop the cached index so the chatbot stops citing deleted content.
    _invalidate_rag_index()

    return {"message": "Material deleted"}
